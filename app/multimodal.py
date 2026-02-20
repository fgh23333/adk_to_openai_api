import re
import base64
import mimetypes
import magic
import io
import csv
from typing import List, Optional, Tuple, Union
import httpx
from app.config import settings
from app.models import ContentPart, ADKPart, ADKInlineData
import logging

logger = logging.getLogger(__name__)


class TextExtractor:
    """提取各种文件格式中的纯文本"""

    @staticmethod
    def extract_from_html(html_content: str) -> str:
        """从 HTML 中提取纯文本"""
        # 移除 script 和 style 标签及其内容
        html_content = re.sub(r'<script[^>]*>.*?</script>', '', html_content, flags=re.DOTALL | re.IGNORECASE)
        html_content = re.sub(r'<style[^>]*>.*?</style>', '', html_content, flags=re.DOTALL | re.IGNORECASE)

        # 移除 HTML 注释
        html_content = re.sub(r'<!--.*?-->', '', html_content, flags=re.DOTALL)

        # 将常见的块级元素替换为换行
        html_content = re.sub(r'<br\s*/?>', '\n', html_content, flags=re.IGNORECASE)
        html_content = re.sub(r'</p>', '\n', html_content, flags=re.IGNORECASE)
        html_content = re.sub(r'</div>', '\n', html_content, flags=re.IGNORECASE)
        html_content = re.sub(r'</li>', '\n', html_content, flags=re.IGNORECASE)
        html_content = re.sub(r'<li[^>]*>', '- ', html_content, flags=re.IGNORECASE)
        html_content = re.sub(r'</h[1-6]>', '\n\n', html_content, flags=re.IGNORECASE)

        # 移除所有剩余的 HTML 标签
        html_content = re.sub(r'<[^>]+>', '', html_content)

        # 解码 HTML 实体
        import html
        html_content = html.unescape(html_content)

        # 清理多余空白
        lines = html_content.split('\n')
        lines = [line.strip() for line in lines]
        text = '\n'.join(line for line in lines if line)

        # 合并多余空行
        text = re.sub(r'\n{3,}', '\n\n', text)

        return text.strip()

    @staticmethod
    def extract_from_markdown(md_content: str) -> str:
        """从 Markdown 中提取纯文本"""
        # 移除代码块
        md_content = re.sub(r'```[\s\S]*?```', '', md_content)
        md_content = re.sub(r'`[^`]+`', '', md_content)

        # 移除链接但保留文本
        md_content = re.sub(r'\[([^\]]+)\]\([^)]+\)', r'\1', md_content)

        # 移除图片
        md_content = re.sub(r'!\[[^\]]*\]\([^)]+\)', '', md_content)

        # 移除标题标记
        md_content = re.sub(r'^#{1,6}\s+', '', md_content, flags=re.MULTILINE)

        # 移除粗体/斜体标记
        md_content = re.sub(r'\*\*([^*]+)\*\*', r'\1', md_content)
        md_content = re.sub(r'\*([^*]+)\*', r'\1', md_content)
        md_content = re.sub(r'__([^_]+)__', r'\1', md_content)
        md_content = re.sub(r'_([^_]+)_', r'\1', md_content)

        # 移除列表标记
        md_content = re.sub(r'^\s*[-*+]\s+', '', md_content, flags=re.MULTILINE)
        md_content = re.sub(r'^\s*\d+\.\s+', '', md_content, flags=re.MULTILINE)

        # 移除引用标记
        md_content = re.sub(r'^\s*>\s?', '', md_content, flags=re.MULTILINE)

        # 移除水平线
        md_content = re.sub(r'^[-*_]{3,}$', '', md_content, flags=re.MULTILINE)

        # 清理多余空白
        lines = md_content.split('\n')
        lines = [line.strip() for line in lines]
        text = '\n'.join(line for line in lines if line)

        return text.strip()

    @staticmethod
    def extract_from_docx(file_data: bytes) -> str:
        """从 DOCX 中提取纯文本"""
        try:
            from docx import Document
            doc = Document(io.BytesIO(file_data))

            text_parts = []
            for para in doc.paragraphs:
                if para.text.strip():
                    text_parts.append(para.text.strip())

            # 提取表格内容
            for table in doc.tables:
                for row in table.rows:
                    row_text = ' | '.join(cell.text.strip() for cell in row.cells if cell.text.strip())
                    if row_text:
                        text_parts.append(row_text)

            return '\n'.join(text_parts)
        except Exception as e:
            logger.error(f"Failed to extract from DOCX: {e}")
            return ""

    @staticmethod
    def extract_from_xlsx(file_data: bytes) -> str:
        """从 XLSX 中提取纯文本"""
        try:
            from openpyxl import load_workbook
            wb = load_workbook(io.BytesIO(file_data), data_only=True)

            text_parts = []
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                text_parts.append(f"[Sheet: {sheet_name}]")

                for row in sheet.iter_rows(values_only=True):
                    row_text = ' | '.join(str(cell) if cell is not None else '' for cell in row)
                    if row_text.strip(' |'):
                        text_parts.append(row_text)

                text_parts.append("")  # 空行分隔不同 sheet

            return '\n'.join(text_parts)
        except Exception as e:
            logger.error(f"Failed to extract from XLSX: {e}")
            return ""

    @staticmethod
    def extract_from_csv(file_data: bytes) -> str:
        """从 CSV 中提取纯文本"""
        try:
            # 尝试检测编码
            content = file_data.decode('utf-8', errors='ignore')

            # 使用 csv 模块解析
            reader = csv.reader(io.StringIO(content))
            rows = list(reader)

            if not rows:
                return ""

            # 格式化输出
            text_parts = []
            for row in rows:
                row_text = ' | '.join(str(cell) for cell in row)
                if row_text.strip(' |'):
                    text_parts.append(row_text)

            return '\n'.join(text_parts)
        except Exception as e:
            logger.error(f"Failed to extract from CSV: {e}")
            # 降级为直接返回原始内容
            return file_data.decode('utf-8', errors='ignore')

    @staticmethod
    def extract_from_pptx(file_data: bytes) -> str:
        """从 PPTX 中提取纯文本"""
        try:
            from pptx import Presentation
            prs = Presentation(io.BytesIO(file_data))

            text_parts = []
            for slide_num, slide in enumerate(prs.slides, 1):
                text_parts.append(f"[Slide {slide_num}]")

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_parts.append(shape.text.strip())

                text_parts.append("")  # 空行分隔不同 slide

            return '\n'.join(text_parts)
        except Exception as e:
            logger.error(f"Failed to extract from PPTX: {e}")
            return ""


class MultimodalProcessor:
    def __init__(self):
        self.max_file_size = settings.max_file_size_mb * 1024 * 1024  # Convert to bytes
        self.timeout = settings.download_timeout
        self.text_extractor = TextExtractor()

        # Gemini/ADK 支持的多模态文件类型 + 需要文本提取的文件类型
        # 参考: https://ai.google.dev/gemini-api/docs/vision
        self.supported_types = {
            "images": [
                "image/jpeg",
                "image/png",
                "image/gif",
                "image/webp",
            ],
            "videos": [
                "video/mp4",
                "video/mpeg",
                "video/quicktime",  # .mov
                "video/x-msvideo",  # .avi
                "video/x-flv",      # .flv
                "video/webm",
                "video/3gpp",       # .3gp
            ],
            "audio": [
                "audio/mpeg",       # .mp3
                "audio/mp3",
                "audio/wav",
                "audio/flac",
                "audio/ogg",
                "audio/aac",
                "audio/mp4",        # .m4a
                "audio/webm",
            ],
            "documents": [
                "application/pdf",
            ],
            "text": [
                "text/plain",
                "text/html",
                "text/markdown",
                "text/x-markdown",
                "text/csv",
            ],
            "office": [
                "application/vnd.openxmlformats-officedocument.wordprocessingml.document",  # .docx
                "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",  # .xlsx
                "application/vnd.openxmlformats-officedocument.presentationml.presentation",  # .pptx
                "application/vnd.ms-excel",  # .xls (legacy)
                "application/msword",  # .doc (legacy)
                "application/vnd.ms-powerpoint",  # .ppt (legacy)
            ]
        }

        # 需要提取文本的文件类型 -> 提取方法
        self.text_extraction_types = {
            "text/html": "html",
            "text/markdown": "markdown",
            "text/x-markdown": "markdown",
            "text/csv": "csv",
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "docx",
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": "xlsx",
            "application/vnd.openxmlformats-officedocument.presentationml.presentation": "pptx",
            "application/vnd.ms-excel": "xls",
            "application/msword": "doc",
            "application/vnd.ms-powerpoint": "ppt",
        }

        # 文件大小限制（字节）
        self.file_size_limits = {
            "image": 20 * 1024 * 1024,      # 20MB
            "video": 200 * 1024 * 1024,     # 200MB
            "audio": 20 * 1024 * 1024,      # 20MB
            "document": 50 * 1024 * 1024,   # 50MB (PDF)
            "text": 10 * 1024 * 1024,       # 10MB
            "office": 50 * 1024 * 1024,     # 50MB
        }
        
    def validate_file(self, file_data: bytes, filename: str, mime_type: str = None) -> Tuple[bool, str, str]:
        """
        验证文件类型和大小
        Returns: (is_valid, error_message, detected_mime_type)
        """
        try:
            # 检测文件大小
            file_size = len(file_data)
            
            # 检测MIME类型
            if not mime_type:
                mime_type = magic.from_buffer(file_data, mime=True)
            
            # 获取所有支持的类型
            all_supported_types = []
            for types in self.supported_types.values():
                all_supported_types.extend(types)
            
            # 检查文件类型
            if mime_type not in all_supported_types:
                return False, f"不支持的文件类型: {mime_type}", mime_type
            
            # 确定文件类别和大小限制
            category = None
            size_limit = self.max_file_size  # 默认限制
            
            for cat, types in self.supported_types.items():
                if mime_type in types:
                    category = cat.rstrip('s')  # 移除复数形式
                    size_limit = self.file_size_limits.get(category, self.max_file_size)
                    break
            
            # 检查文件大小
            if file_size > size_limit:
                size_mb = size_limit / (1024 * 1024)
                return False, f"文件大小超过限制 ({size_mb:.1f}MB)", mime_type
            
            return True, "", mime_type
            
        except Exception as e:
            logger.error(f"文件验证失败: {e}")
            return False, f"文件验证失败: {str(e)}", ""

    def process_base64_file(self, base64_data: str, filename: str, mime_type: str = None) -> Tuple[Optional[ADKInlineData], Optional[str]]:
        """
        处理Base64编码的文件数据
        Returns: (inline_data, extracted_text) - inline_data 用于二进制文件，extracted_text 用于需要文本提取的文件
        """
        try:
            # 解码Base64数据
            if base64_data.startswith(f"data:"):
                # 移除数据URL前缀
                base64_data = base64_data.split(",", 1)[1] if "," in base64_data else base64_data

            file_data = base64.b64decode(base64_data)

            # 验证文件
            is_valid, error_msg, detected_mime = self.validate_file(file_data, filename, mime_type)

            if not is_valid:
                logger.error(f"文件验证失败: {error_msg}")
                return None, None

            # 检查是否需要提取文本
            if detected_mime in self.text_extraction_types:
                extraction_type = self.text_extraction_types[detected_mime]
                extracted_text = ""

                if extraction_type == "html":
                    text_content = file_data.decode('utf-8', errors='ignore')
                    extracted_text = self.text_extractor.extract_from_html(text_content)
                    logger.info(f"Extracted text from HTML: {len(extracted_text)} chars")
                elif extraction_type == "markdown":
                    text_content = file_data.decode('utf-8', errors='ignore')
                    extracted_text = self.text_extractor.extract_from_markdown(text_content)
                    logger.info(f"Extracted text from Markdown: {len(extracted_text)} chars")
                elif extraction_type == "csv":
                    extracted_text = self.text_extractor.extract_from_csv(file_data)
                    logger.info(f"Extracted text from CSV: {len(extracted_text)} chars")
                elif extraction_type == "docx":
                    extracted_text = self.text_extractor.extract_from_docx(file_data)
                    logger.info(f"Extracted text from DOCX: {len(extracted_text)} chars")
                elif extraction_type == "xlsx":
                    extracted_text = self.text_extractor.extract_from_xlsx(file_data)
                    logger.info(f"Extracted text from XLSX: {len(extracted_text)} chars")
                elif extraction_type == "pptx":
                    extracted_text = self.text_extractor.extract_from_pptx(file_data)
                    logger.info(f"Extracted text from PPTX: {len(extracted_text)} chars")
                elif extraction_type in ("doc", "xls", "ppt"):
                    # 旧版 Office 格式暂时不支持，返回提示
                    extracted_text = f"[Legacy Office format ({extraction_type}) - please convert to modern format]"
                    logger.warning(f"Legacy Office format not fully supported: {extraction_type}")
                else:
                    # 尝试作为文本解码
                    extracted_text = file_data.decode('utf-8', errors='ignore')

                return None, extracted_text

            # 对于其他类型，返回 inline data
            final_base64 = base64.b64encode(file_data).decode('utf-8')

            return ADKInlineData(
                mimeType=detected_mime,
                data=final_base64
            ), None

        except Exception as e:
            logger.error(f"处理Base64文件失败: {e}")
            return None, None

    async def process_content(self, content_parts: List[ContentPart]) -> Tuple[str, List[ADKPart]]:
        """
        Process content parts, extracting text and handling multimodal content.
        Supports: text, image_url, audio_url, video_url, input_audio, file
        Returns tuple of (combined_text, adk_parts)
        """
        text_parts = []
        adk_parts = []

        logger.info(f"Starting multimodal processing for {len(content_parts)} content parts")

        for i, part in enumerate(content_parts):
            logger.info(f"Processing content part {i}: type={part.type}")

            if part.type == "text" and part.text:
                logger.info(f"Found text part: {part.text[:100]}...")
                text_parts.append(part.text)
                # Extract URLs from text for video/file processing
                urls = self._extract_urls_from_text(part.text)
                logger.info(f"Found {len(urls)} URLs in text: {urls}")
                for url in urls:
                    try:
                        logger.info(f"Attempting to download URL: {url}")
                        inline_data = await self._download_and_convert_url(url)
                        if inline_data:
                            logger.info(f"Successfully downloaded and converted URL: {inline_data.mimeType}")
                            adk_parts.append(ADKPart(inlineData=inline_data))
                        else:
                            logger.warning(f"Failed to download URL: {url} - no data returned")
                    except Exception as e:
                        logger.error(f"Failed to process URL {url}: {e}")

            elif part.type == "image_url" and part.image_url:
                await self._process_image_url(part.image_url.url, adk_parts)

            elif part.type == "audio_url" and part.audio_url:
                await self._process_audio_url(part.audio_url.url, adk_parts)

            elif part.type == "video_url" and part.video_url:
                await self._process_video_url(part.video_url.url, adk_parts)

            elif part.type == "input_audio" and part.input_audio:
                self._process_input_audio(part.input_audio, adk_parts)

            elif part.type == "file" and part.file:
                await self._process_file(part.file, adk_parts)

            else:
                logger.warning(f"Unsupported content part type: {part.type}")

        # Combine all text parts
        combined_text = " ".join(text_parts)

        # Add text part if we have combined text
        if combined_text.strip():
            adk_parts.insert(0, ADKPart(text=combined_text))

        return combined_text, adk_parts

    async def _process_image_url(self, url: str, adk_parts: List[ADKPart]):
        """处理图片URL（支持URL和Base64数据）"""
        logger.info(f"Found image_url part: {url[:100]}...")

        if url.startswith("data:"):
            try:
                logger.info(f"Processing Base64 image data")
                mime_type = None
                if ":" in url:
                    mime_type = url.split(":")[1].split(";")[0]

                inline_data, extracted_text = self.process_base64_file(url, "image", mime_type)
                if inline_data:
                    logger.info(f"Successfully processed Base64 image: {inline_data.mimeType}")
                    adk_parts.append(ADKPart(inlineData=inline_data))
                elif extracted_text:
                    # 对于图片URL来说不太可能有文本提取，但以防万一
                    adk_parts.append(ADKPart(text=f"[Image content]\n{extracted_text}"))
                else:
                    logger.warning(f"Failed to process Base64 image data")
            except Exception as e:
                logger.error(f"Failed to process Base64 image: {e}")
        else:
            try:
                logger.info(f"Attempting to download image URL: {url}")
                inline_data = await self._download_and_convert_url(url)
                if inline_data:
                    logger.info(f"Successfully downloaded and converted image: {inline_data.mimeType}")
                    adk_parts.append(ADKPart(inlineData=inline_data))
                else:
                    logger.warning(f"Failed to download image URL: {url}")
            except Exception as e:
                logger.error(f"Failed to process image URL {url}: {e}")

    async def _process_audio_url(self, url: str, adk_parts: List[ADKPart]):
        """处理音频URL"""
        logger.info(f"Found audio_url part: {url[:100]}...")

        if url.startswith("data:"):
            try:
                logger.info(f"Processing Base64 audio data")
                mime_type = None
                if ":" in url:
                    mime_type = url.split(":")[1].split(";")[0]

                inline_data, extracted_text = self.process_base64_file(url, "audio", mime_type)
                if inline_data:
                    logger.info(f"Successfully processed Base64 audio: {inline_data.mimeType}")
                    adk_parts.append(ADKPart(inlineData=inline_data))
                else:
                    logger.warning(f"Failed to process Base64 audio data")
            except Exception as e:
                logger.error(f"Failed to process Base64 audio: {e}")
        else:
            try:
                logger.info(f"Attempting to download audio URL: {url}")
                inline_data = await self._download_and_convert_url(url)
                if inline_data:
                    logger.info(f"Successfully downloaded and converted audio: {inline_data.mimeType}")
                    adk_parts.append(ADKPart(inlineData=inline_data))
                else:
                    logger.warning(f"Failed to download audio URL: {url}")
            except Exception as e:
                logger.error(f"Failed to process audio URL {url}: {e}")

    async def _process_video_url(self, url: str, adk_parts: List[ADKPart]):
        """处理视频URL"""
        logger.info(f"Found video_url part: {url[:100]}...")

        if url.startswith("data:"):
            try:
                logger.info(f"Processing Base64 video data")
                mime_type = None
                if ":" in url:
                    mime_type = url.split(":")[1].split(";")[0]

                inline_data, extracted_text = self.process_base64_file(url, "video", mime_type)
                if inline_data:
                    logger.info(f"Successfully processed Base64 video: {inline_data.mimeType}")
                    adk_parts.append(ADKPart(inlineData=inline_data))
                else:
                    logger.warning(f"Failed to process Base64 video data")
            except Exception as e:
                logger.error(f"Failed to process Base64 video: {e}")
        else:
            try:
                logger.info(f"Attempting to download video URL: {url}")
                inline_data = await self._download_and_convert_url(url)
                if inline_data:
                    logger.info(f"Successfully downloaded and converted video: {inline_data.mimeType}")
                    adk_parts.append(ADKPart(inlineData=inline_data))
                else:
                    logger.warning(f"Failed to download video URL: {url}")
            except Exception as e:
                logger.error(f"Failed to process video URL {url}: {e}")

    def _process_input_audio(self, input_audio, adk_parts: List[ADKPart]):
        """处理 input_audio 类型（OpenAI格式的音频输入）"""
        logger.info(f"Found input_audio part, format: {input_audio.format}")

        try:
            # 根据format推断MIME类型
            format_to_mime = {
                "mp3": "audio/mpeg",
                "wav": "audio/wav",
                "flac": "audio/flac",
                "aac": "audio/aac",
                "ogg": "audio/ogg",
                "m4a": "audio/mp4",
                "webm": "audio/webm"
            }
            mime_type = format_to_mime.get(input_audio.format.lower(), f"audio/{input_audio.format}")

            # 创建内联数据
            inline_data = ADKInlineData(
                mimeType=mime_type,
                data=input_audio.data
            )
            logger.info(f"Successfully processed input_audio: {mime_type}")
            adk_parts.append(ADKPart(inlineData=inline_data))
        except Exception as e:
            logger.error(f"Failed to process input_audio: {e}")

    async def _process_file(self, file_content, adk_parts: List[ADKPart]):
        """处理通用文件类型"""
        logger.info(f"Found file part, filename: {file_content.filename}")

        try:
            # 优先使用Base64数据
            if file_content.data:
                mime_type = file_content.mime_type or "application/octet-stream"

                # 使用 process_base64_file 处理（它会处理文本提取）
                inline_data, extracted_text = self.process_base64_file(
                    file_content.data,
                    file_content.filename or "file",
                    mime_type
                )

                if inline_data:
                    logger.info(f"Successfully processed file data: {inline_data.mimeType}")
                    adk_parts.append(ADKPart(inlineData=inline_data))
                elif extracted_text:
                    # 文本提取后的内容作为文本发送
                    logger.info(f"Extracted text from file: {len(extracted_text)} chars")
                    adk_parts.append(ADKPart(text=f"[File: {file_content.filename}]\n{extracted_text}"))
                else:
                    logger.warning(f"Failed to process file data")

            # 如果没有数据但有URL，尝试下载
            elif file_content.url:
                logger.info(f"Attempting to download file URL: {file_content.url}")
                inline_data = await self._download_and_convert_url(file_content.url)
                if inline_data:
                    # 检查是否需要文本提取
                    if inline_data.mimeType in self.text_extraction_types:
                        file_data = base64.b64decode(inline_data.data)
                        text_content = file_data.decode('utf-8', errors='ignore')
                        extraction_type = self.text_extraction_types[inline_data.mimeType]

                        if extraction_type == "html":
                            extracted_text = self.text_extractor.extract_from_html(text_content)
                        elif extraction_type == "markdown":
                            extracted_text = self.text_extractor.extract_from_markdown(text_content)
                        else:
                            extracted_text = text_content

                        logger.info(f"Extracted text from downloaded file: {len(extracted_text)} chars")
                        adk_parts.append(ADKPart(text=f"[File from URL]\n{extracted_text}"))
                    else:
                        logger.info(f"Successfully downloaded and converted file: {inline_data.mimeType}")
                        adk_parts.append(ADKPart(inlineData=inline_data))
                else:
                    logger.warning(f"Failed to download file URL: {file_content.url}")
            else:
                logger.warning("File content has neither data nor url")

        except Exception as e:
            logger.error(f"Failed to process file: {e}")
    
    def _extract_urls_from_text(self, text: str) -> List[str]:
        """Extract HTTP/HTTPS URLs and file paths from text using regex."""
        urls = []
        
        # Extract HTTP/HTTPS URLs
        http_pattern = r'https?://[^\s<>"{}|\\^`\[\]]+'
        http_urls = re.findall(http_pattern, text)
        urls.extend(http_urls)
        
        # Extract file:// URLs
        file_pattern = r'file://[^\s<>"{}|\\^`\[\]]+'
        file_urls = re.findall(file_pattern, text)
        urls.extend(file_urls)
        
        # Extract Windows file paths (e.g., D:\folder\file.txt)
        windows_pattern = r'[A-Za-z]:\\[^\s<>"{}|\\^`\[\]]+\.[A-Za-z0-9]+'
        windows_paths = re.findall(windows_pattern, text)
        urls.extend(windows_paths)
        
        logger.info(f"Extracted URLs: {urls}")
        return urls
    
    async def _download_and_convert_url(self, url: str) -> Optional[ADKInlineData]:
        """
        Download file from URL and convert to base64 inline data.
        Returns None if download fails or file is too large.
        """
        logger.info(f"Starting download of URL: {url}")
        try:
            async with httpx.AsyncClient(timeout=self.timeout) as client:
                # First, get content type and file size with HEAD request
                logger.info(f"Sending HEAD request to: {url}")
                head_response = await client.head(url)
                logger.info(f"HEAD response status: {head_response.status_code}")
                
                content_type = head_response.headers.get('content-type', '')
                content_length = head_response.headers.get('content-length')
                
                logger.info(f"Content-Type: {content_type}, Content-Length: {content_length}")
                
                # Check file size
                if content_length and int(content_length) > self.max_file_size:
                    logger.warning(f"File too large: {content_length} bytes > {self.max_file_size} bytes")
                    return None
                
                # Download the file
                logger.info(f"Starting GET request to download file")
                response = await client.get(url)
                logger.info(f"GET response status: {response.status_code}")
                response.raise_for_status()
                
                # Check actual file size
                actual_size = len(response.content)
                logger.info(f"Downloaded {actual_size} bytes")
                
                if actual_size > self.max_file_size:
                    logger.warning(f"Downloaded file too large: {actual_size} bytes")
                    return None
                
                # Determine MIME type
                if not content_type:
                    content_type, _ = mimetypes.guess_type(url)
                    if not content_type:
                        # Default to binary if we can't determine the type
                        content_type = 'application/octet-stream'
                
                logger.info(f"Final MIME type: {content_type}")
                
                # Convert to base64
                base64_data = base64.b64encode(response.content).decode('utf-8')
                logger.info(f"Converted to base64, length: {len(base64_data)}")
                
                return ADKInlineData(
                    mimeType=content_type,
                    data=base64_data
                )
                
        except httpx.TimeoutException:
            logger.error(f"Timeout downloading URL: {url}")
            return None
        except httpx.HTTPStatusError as e:
            logger.error(f"HTTP error downloading URL {url}: {e.response.status_code}")
            return None
        except Exception as e:
            logger.error(f"Error downloading URL {url}: {e}")
            return None